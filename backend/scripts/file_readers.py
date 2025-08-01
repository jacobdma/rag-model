# Standard imports
import threading
import time
from io import BytesIO, StringIO
from pathlib import Path

# Third-party imports
import docx2txt
import fitz
import pandas as pd
from pptx import Presentation

_pdf_lock = threading.Lock()

def read_docx(f):
    if isinstance(f, Path):
        text = docx2txt.process(str(f))
    else:
        text = docx2txt.process(BytesIO(f))
    return text

def read_pptx(f):
    if isinstance(f, (bytes, bytearray)):
        f = BytesIO(f)
    presentation = Presentation(f)
    text_lines = [
        paragraph.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for text_frame in [getattr(shape, "text_frame", None)]
        if text_frame is not None
        for paragraph in text_frame.paragraphs
    ]
    return "\n".join(text_lines).strip()

def read_txt(f):
    if isinstance(f, (bytes, bytearray)):
        return f.decode("utf-8", errors="ignore")
    return f.read().decode("utf-8", errors="ignore")

def read_pdf(f, filename):
    start_pdf = time.time()
    if hasattr(f, "seek"):
        f.seek(0)
    with _pdf_lock:
        with fitz.open(stream=f.read(), filetype="pdf") as doc:
            text_parts = [page.get_text() for page in doc]  # type: ignore
    if time.time() - start_pdf > 900:
        return None
    return "\n".join(text_parts)

def read_csv(f):
    try:
        if isinstance(f, (bytes, bytearray)):
            f = StringIO(f.decode("utf-8"))
        return pd.read_csv(f).to_string(index=False)
    except Exception:
        # Wrap in StringIO
        if not isinstance(f, StringIO):
            if isinstance(f, (bytes, bytearray)):
                f = StringIO(f.decode("latin1"))
            else:
                f = StringIO(str(f))
        f.seek(0)
        return pd.read_csv(f, encoding="latin1").to_string(index=False)

class FileReader:
    def __init__(self, _supported_exts):
        self._supported_exts = _supported_exts
        
    def read_docs(self, file: Path):
        readers = {
            ".docx": read_docx,
            ".pptx": read_pptx,
            ".txt": read_txt,
            ".pdf": lambda b: read_pdf(BytesIO(b), filename),
            ".csv": read_csv,
        }
        filename = str(file)
        ext = file.suffix.lower()
        if ext not in self._supported_exts:
            return None
        with open(file, "rb") as f:
            file_bytes = f.read()
        reader = readers.get(ext)
        if reader is None:
            return None, filename
        return reader(file_bytes), filename